from pptx import Presentation

# Create a presentation object
prs = Presentation()

# Define a function to add a slide with title and content
def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]  # Use the title and content layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    content_placeholder.text = content

# Slide 1: Introduction
add_slide("Introduction to AI and Machine Learning", "")

# Slide 2: Introduction
add_slide("Introduction",
          "Personal background in software engineering and AI/ML.\n"
          "Motivation for teaching AI and ML.")

# Slide 3: What is AI and Why We Need to Learn It
add_slide("What is AI and Why We Need to Learn It",
          "Definition of AI:\n"
          "- AI is the simulation of human intelligence in machines.\n"
          "Importance:\n"
          "- Enhancing efficiency and productivity.\n"
          "- Solving complex problems.\n"
          "- Driving innovation in various fields.\n"
          "Applications of AI:\n"
          "- Healthcare, finance, transportation, etc.")

# Slide 4: Types and Branches of AI
add_slide("Types and Branches of AI",
          "Types of AI:\n"
          "- Narrow AI: AI designed for specific tasks.\n"
          "- General AI: AI with generalized human cognitive abilities (future goal).\n"
          "- Super AI: AI that surpasses human intelligence (theoretical concept).\n"
          "Branches of AI:\n"
          "- Machine Learning\n"
          "- Natural Language Processing\n"
          "- Computer Vision\n"
          "- Robotics\n"
          "- Expert Systems")

# Slide 5: Applications of Each Branch
add_slide("Applications of Each Branch",
          "Machine Learning:\n"
          "- Predictive analytics, recommendation systems.\n"
          "Natural Language Processing:\n"
          "- Chatbots, language translation.\n"
          "Computer Vision:\n"
          "- Image and video analysis, facial recognition.\n"
          "Robotics:\n"
          "- Autonomous vehicles, industrial robots.\n"
          "Expert Systems:\n"
          "- Decision support systems, medical diagnosis.")

# Slide 6: Real-Life Applications of AI
add_slide("Real-Life Applications of AI",
          "Healthcare:\n"
          "- Predictive diagnostics, personalized treatment plans.\n"
          "Finance:\n"
          "- Fraud detection, algorithmic trading.\n"
          "Transportation:\n"
          "- Autonomous driving, traffic management.\n"
          "Retail:\n"
          "- Customer service chatbots, inventory management.\n"
          "Entertainment:\n"
          "- Content recommendation, video game AI.")

# Slide 7: Roadmap to Study AI and Resources
add_slide("Roadmap to Study AI and Resources",
          "Foundations:\n"
          "- Mathematics (linear algebra, calculus, probability).\n"
          "- Programming (Python, R).\n"
          "Core AI Concepts:\n"
          "- Algorithms and data structures.\n"
          "- Basics of machine learning and deep learning.\n"
          "Specialization:\n"
          "- Choose a branch of AI to specialize in.\n"
          "Resources:\n"
          "- Books: 'Artificial Intelligence: A Modern Approach' by Stuart Russell and Peter Norvig.\n"
          "- Online Courses: Coursera, edX, Udacity.\n"
          "- Tutorials and Documentation: TensorFlow, PyTorch, Scikit-Learn.")

# Slide 8: Introduction to Machine Learning and Its Types
add_slide("Introduction to Machine Learning and Its Types",
          "Definition of Machine Learning:\n"
          "- ML is the ability of machines to learn from data.\n"
          "Types of Machine Learning:\n"
          "- Supervised Learning: Learning from labeled data.\n"
          "- Unsupervised Learning: Learning from unlabeled data.\n"
          "- Reinforcement Learning: Learning through trial and error.")

# Slide 9: Types of Supervised Machine Learning
add_slide("Types of Supervised Machine Learning",
          "Types of Supervised Learning:\n"
          "- Regression: Predicting continuous outcomes.\n"
          "- Classification: Predicting categorical outcomes.")

# Slide 10: Detailed Study of Regression Algorithms
add_slide("Detailed Study of Regression Algorithms",
          "Introduction to Regression:\n"
          "- Method for predicting continuous outcomes.\n"
          "Types of Regression:\n"
          "- Linear Regression:\n"
          "  - Equation: y = mx + c\n"
          "  - Applications: Predicting house prices, stock prices.\n"
          "  - Algorithms: Ordinary Least Squares (OLS), Gradient Descent.\n"
          "- Multiple Linear Regression:\n"
          "  - Equation: y = b_0 + b_1x_1 + b_2x_2 + ... + b_nx_n\n"
          "  - Applications: Predicting outcomes based on multiple factors.\n"
          "  - Algorithms: Ordinary Least Squares (OLS), Gradient Descent.\n"
          "- Polynomial Regression:\n"
          "  - Equation: y = b_0 + b_1x + b_2x^2 + ... + b_nx^n\n"
          "  - Applications: Modeling non-linear relationships.\n"
          "  - Algorithms: Ordinary Least Squares (OLS), Gradient Descent.\n"
          "- Ridge Regression:\n"
          "  - Concept: Linear regression with L2 regularization.\n"
          "  - Applications: When there are multicollinearity issues.\n"
          "  - Algorithms: Closed-form solution, Gradient Descent.\n"
          "- Lasso Regression:\n"
          "  - Concept: Linear regression with L1 regularization.\n"
          "  - Applications: Feature selection, sparse models.\n"
          "  - Algorithms: Coordinate Descent.\n"
          "- Elastic Net Regression:\n"
          "  - Concept: Combines L1 and L2 regularization.\n"
          "  - Applications: Feature selection and multicollinearity.\n"
          "  - Algorithms: Coordinate Descent.\n"
          "Evaluation Metrics:\n"
          "- R-squared, Adjusted R-squared, Mean Squared Error (MSE), Root Mean Squared Error (RMSE).\n"
          "Practical Considerations:\n"
          "- Overfitting and underfitting.\n"
          "- Model validation and cross-validation.")

# Slide 11: Practical Application Using Linear Regression
add_slide("Practical Application Using Linear Regression",
          "Dataset Selection:\n"
          "- Choose a simple dataset (e.g., housing prices).\n"
          "Implementation Steps:\n"
          "- Data Loading and Preprocessing:\n"
          "  - Load dataset using pandas.\n"
          "  - Handle missing values, normalize features.\n"
          "- Model Training:\n"
          "  - Split data into training and testing sets.\n"
          "  - Train a linear regression model using Scikit-Learn.\n"
          "- Model Evaluation:\n"
          "  - Make predictions on test data.\n"
          "  - Evaluate model performance using discussed metrics.\n"
          "- Visualization:\n"
          "  - Plot the regression line and data points.")

# Slide 12: Q&A and Discussion
add_slide("Q&A and Discussion",
          "Encourage questions and discussion on the presented topics.\n"
          "Address common challenges and best practices in AI and ML.")

# Save the presentation
prs.save('AI_and_ML_Presentation.pptx')
